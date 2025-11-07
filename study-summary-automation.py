#the purpose of this automation tool is to summarize notes for students who wants to maximize their time in studying
#The tool will get the input as .pdf or .ppt and will return .docx with highlights and summary for better cognitive presentation

import sys
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
from transformers import BartTokenizer, BartForConditionalGeneration, pipeline

#************************************************************************
# For checking if for the location of local cache memory for Bart
#************************************************************************
#from huggingface_hub import hf_hub_download, constants
#print(constants.HF_HUB_CACHE)



#*************************************************************************
# EXTRACTION OF DOCUMENT
#*************************************************************************
#Function that converts pdf into text format
def pdfToText(pdf_path):
    openPDFFile = fitz.open(pdf_path)
    text = ""
    #store all the text inside the pdf
    for page in openPDFFile:
        text += page.get_text()
    return text

#Function that converts ppt presentation into text format
def pptToText(ppt_path):

     #to be follow for the pptx presentation - status buggy
    #*******************************************************
    openPPTFile = Presentation(ppt_path)
    text = ""
    #store all the text inside the ppt including text inside shape
    for slide in openPPTFile.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

#Function that will write the extracted information into .docx format
def createDocxFile(context, filename):
        createdDoc = Document()
        createdDoc.add_paragraph(context)
        createdDoc.save(filename)
        print(f"Saved: {filename}")

#**********************************************************************
# Summarize text for note taking and filtered information
#**********************************************************************
#Fuction that will read and extract all text from the document 
def readDocument(filename):
        readDocument = Document(filename)
        return "\n".join([paragraph.text for paragraph in readDocument.paragraphs])

#Function that will summarize context 
def summarizeContext(context):
        print("Initializing summarizer... it may take a while please wait...")
        
        # Load from cache only
        tokenizer = BartTokenizer.from_pretrained("facebook/bart-large-cnn", local_files_only=True)
        model = BartForConditionalGeneration.from_pretrained("facebook/bart-large-cnn", local_files_only=True)

        #Using bart model to perform the summarizer
        summarizer = pipeline(
             "summarization",
              model=model,
              tokenizer=tokenizer)

        chunks = []
        # Limitation to prevent overflowing
        chunkSize = 900


        for i in range(0, len(context), chunkSize):
             chunks.append(context[i:i+ chunkSize])

        summaries = []

        for i, chunk in enumerate(chunks):
             print(f" > Summarizing chunk {i+1}/{len(chunks)}")
             
             #AI text-summarization pipeline for model (BART)
             result = summarizer(chunk, max_length=200, min_length=80, do_sample=False)[0]['summary_text']
             summaries.append(result)
        
        return "\n\n".join(summaries)



def main():
    if len(sys.argv) < 3:
        #Instruct user on the format of using the tool.
        print("\nUsage:")
        print("  python3 study-summary-automation.py extract <input.pdf|input.pptx>")
        print("  python3 study-summary-automation.py summarize <input.docx>\n")
        sys.exit(1)

    instruction = sys.argv[1]
    documentPath = sys.argv[2]


    #choosing instruction 
    match instruction:
        case "extract":
            print("\nExtracting data from a document(.pdf)... Please wait ...")

            #validates if file is either pdf or ppt
            if documentPath.endswith(".pdf"):
                context = pdfToText(documentPath)

            #to be follow for the pptx presentation
            #********************************************************
            #elif documentPath.endswith(".pptx"):
                #context = pptToText(documentPath)
            else:
                print("Unsupported file format. Kindly use PDF to use the tool.")
                sys.exit(1)

            #paste all the extracted data into .docx format
            createDocxFile(context, f"{documentPath}_extracted.docx")
        
        case "summarize":
             print("\nConverting document(.docx) into study notes... Please wait...")

             context = readDocument(documentPath)
             summarized = summarizeContext(context)
             createDocxFile(summarized, f"{documentPath}_summary.docx")

        case _: #default case - if neither of the two selection above
              print("Invalid command. Use either 'extract' or 'summarize'." )


if __name__ == "__main__":
    main()
